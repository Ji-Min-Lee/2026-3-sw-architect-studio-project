from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# Colors
BG = RGBColor(0x0D, 0x11, 0x17)
WHITE = RGBColor(0xE6, 0xED, 0xF3)
BLUE = RGBColor(0x58, 0xA6, 0xFF)
GREEN = RGBColor(0x3F, 0xB9, 0x50)
RED = RGBColor(0xF8, 0x51, 0x49)
ORANGE = RGBColor(0xF0, 0x88, 0x3E)
GRAY = RGBColor(0x8B, 0x94, 0x9E)
DARK_CARD = RGBColor(0x16, 0x1B, 0x22)
BLUE_DARK = RGBColor(0x0D, 0x2A, 0x4A)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

def blank_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    return slide

def add_rect(slide, x, y, w, h, fill=None, line_color=None, line_width=Pt(1)):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h, size=Pt(18), bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    run.font.name = "Calibri"
    return txBox

def section_label(slide, text):
    add_text(slide, text, Inches(0.4), Inches(0.15), Inches(4), Inches(0.3),
             size=Pt(11), color=GRAY)

def slide_title(slide, text, y=Inches(0.45)):
    add_text(slide, text, Inches(0.4), y, Inches(12.5), Inches(0.6),
             size=Pt(32), bold=True, color=WHITE)

def footer(slide, page_num):
    add_text(slide, "Blue Sky (Team 3)", Inches(0.4), Inches(7.1), Inches(4), Inches(0.3),
             size=Pt(11), color=GRAY)
    add_text(slide, str(page_num), Inches(12.5), Inches(7.1), Inches(0.5), Inches(0.3),
             size=Pt(11), color=GRAY, align=PP_ALIGN.RIGHT)

# ── SLIDE 1: Title ──────────────────────────────────────────────
s = blank_slide(prs)
add_text(s, "TimeGrapher", Inches(0), Inches(2.2), W, Inches(1.0),
         size=Pt(54), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "From Tick to Trace: Real-Time Acoustic Analysis", Inches(0), Inches(3.3), W, Inches(0.6),
         size=Pt(22), color=BLUE, align=PP_ALIGN.CENTER)
add_text(s, "Final Demonstration — Milestone 3  |  July 1, 2026", Inches(0), Inches(4.0), W, Inches(0.4),
         size=Pt(16), color=GRAY, align=PP_ALIGN.CENTER)
add_text(s, "Blue Sky · Team 3  |  LG SW Architect Training Program × CMU MSE", Inches(0), Inches(4.5), W, Inches(0.4),
         size=Pt(14), color=GRAY, align=PP_ALIGN.CENTER)
# thin line
add_rect(s, Inches(1.5), Inches(5.1), Inches(10.3), Inches(0.02), fill=GRAY)

# ── SLIDE 2: Agenda ──────────────────────────────────────────────
s = blank_slide(prs)
footer(s, 2)
slide_title(s, "Agenda")
items = [
    ("01", "QA Requirements & Tradeoffs"),
    ("02", "Architecture"),
    ("03", "Experiments & Evaluation"),
    ("04", "Use of AI"),
    ("05", "Lessons Learned"),
]
for i, (num, label) in enumerate(items):
    y = Inches(1.3 + i * 0.9)
    add_rect(s, Inches(0.4), y + Inches(0.05), Inches(0.5), Inches(0.5), fill=BLUE)
    add_text(s, num, Inches(0.4), y, Inches(0.5), Inches(0.6),
             size=Pt(16), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, label, Inches(1.1), y, Inches(10), Inches(0.6),
             size=Pt(20), bold=True, color=WHITE)
add_text(s, "Graphs · GUI enhancements · Bonus features  →  demonstrated live on Raspberry Pi",
         Inches(0.4), Inches(6.9), Inches(12), Inches(0.4),
         size=Pt(13), color=GRAY, italic=True)

# ── SLIDE 3: QA Priority Map ──────────────────────────────────────────────
s = blank_slide(prs)
footer(s, 3)
section_label(s, "01 / QA Requirements")
slide_title(s, "Quality Attribute Priorities")

rows = [
    ("1st — GOVERNING", "Measurement Accuracy (QAS-1)", "The reason the system exists. Every other QA serves this one.", BLUE, WHITE),
    ("2nd", "Real-Time Performance (QAS-2)", "21ms audio deadline per block — must never be missed", DARK_CARD, WHITE),
    ("3rd", "Low Latency (QAS-3)", "Beat capture → GUI update < 100ms end-to-end", DARK_CARD, WHITE),
    ("4th", "Extensibility (QAS-4)", "Add new graph tab in ≤ 3 files changed", DARK_CARD, WHITE),
    ("5th", "Correctness (QAS-5)", "False trigger rate < 1%, T1 detection rate > 99%", DARK_CARD, WHITE),
]
col_w = [Inches(1.8), Inches(3.2), Inches(7.0)]
col_x = [Inches(0.3), Inches(2.2), Inches(5.5)]
row_h = Inches(0.82)
start_y = Inches(1.2)

# Header
for ci, hdr in enumerate(["Priority", "Quality Attribute", "Why It Matters"]):
    add_rect(s, col_x[ci], start_y, col_w[ci], Inches(0.4), fill=RGBColor(0x21, 0x26, 0x2E))
    add_text(s, hdr, col_x[ci] + Inches(0.05), start_y, col_w[ci], Inches(0.4),
             size=Pt(13), bold=True, color=GRAY)

for ri, (p, qa, why, bg, tc) in enumerate(rows):
    y = start_y + Inches(0.4) + ri * row_h
    for ci, txt in enumerate([p, qa, why]):
        add_rect(s, col_x[ci], y, col_w[ci], row_h - Inches(0.05), fill=bg)
        add_text(s, txt, col_x[ci] + Inches(0.08), y + Inches(0.05),
                 col_w[ci] - Inches(0.1), row_h - Inches(0.1),
                 size=Pt(13), color=tc, bold=(ci == 0))

add_text(s, "Scope: 28,800 BPH mechanical watch measured against WeiShi No.1000 reference",
         Inches(0.3), Inches(7.0), Inches(12), Inches(0.35), size=Pt(12), color=GRAY, italic=True)

# ── SLIDE 4: QA Tradeoffs ──────────────────────────────────────────────
s = blank_slide(prs)
footer(s, 4)
section_label(s, "01 / QA Requirements")
slide_title(s, "Key Architectural Tradeoffs")

cards = [
    ("TP-1: 96kHz Sample Rate", "↑ Beat Error resolution (QAS-1)", "↑ CPU load risk (QAS-2)", "EXP-02: 0 dropped blocks ✓"),
    ("TP-2: Ring Buffer", "↑ DSP / GUI decoupled (QAS-2, 3)", "+21ms propagation delay", "EXP-03: E2E 2.05ms ✓"),
    ("TP-3: Lazy Rendering", "↑ 85% fewer render calls (QAS-2)", "Non-visible tabs skip updates", "Tab catches up on show() ✓"),
]
card_w = Inches(3.9)
card_h = Inches(3.5)
gap = Inches(0.2)
start_x = Inches(0.35)
card_y = Inches(1.2)

for i, (title, pro, con, res) in enumerate(cards):
    cx = start_x + i * (card_w + gap)
    add_rect(s, cx, card_y, card_w, card_h, fill=DARK_CARD, line_color=BLUE, line_width=Pt(1.5))
    add_text(s, title, cx + Inches(0.1), card_y + Inches(0.1), card_w - Inches(0.2), Inches(0.5),
             size=Pt(14), bold=True, color=BLUE)
    add_text(s, pro, cx + Inches(0.1), card_y + Inches(0.7), card_w - Inches(0.2), Inches(0.5),
             size=Pt(13), color=GREEN)
    add_text(s, con, cx + Inches(0.1), card_y + Inches(1.25), card_w - Inches(0.2), Inches(0.5),
             size=Pt(13), color=RED)
    add_rect(s, cx + Inches(0.1), card_y + Inches(2.5), card_w - Inches(0.2), Inches(0.02), fill=GRAY)
    add_text(s, res, cx + Inches(0.1), card_y + Inches(2.6), card_w - Inches(0.2), Inches(0.6),
             size=Pt(13), color=GREEN, bold=True)

# Bottom row
add_rect(s, Inches(0.35), Inches(5.0), Inches(6.0), Inches(1.0), fill=RGBColor(0x0D, 0x2A, 0x17))
add_text(s, "✓ Deadline miss: 43% → 0%   ✓ E2E latency: 2.05ms   ✓ 0 dropped blocks at 96kHz",
         Inches(0.5), Inches(5.05), Inches(5.7), Inches(0.9), size=Pt(13), color=GREEN)
add_rect(s, Inches(6.6), Inches(5.0), Inches(6.4), Inches(1.0), fill=RGBColor(0x2A, 0x1A, 0x0D))
add_text(s, "⚠ WeiShi accuracy: validated W5\n⚠ Full 14-tab rendering: conditional on EXP-05",
         Inches(6.75), Inches(5.05), Inches(6.1), Inches(0.9), size=Pt(13), color=ORANGE)

# ── SLIDE 5: Architecture Overview ──────────────────────────────────────────────
s = blank_slide(prs)
footer(s, 5)
section_label(s, "02 / Architecture")
slide_title(s, "4-Layer Architecture")

layers = [
    (BLUE,   "PRESENTATION", "Qt Widgets · Graph Tabs · UI Controls"),
    (GREEN,  "DOMAIN",       "MeasurementEngine · Value Objects · Observer"),
    (ORANGE, "ENGINE",       "DSP Pipeline · BeatDetector · FilterChain"),
    (GRAY,   "INFRASTRUCTURE","AudioCapture · IAudioSource · ALSA / Qt Multimedia"),
]
lx = Inches(0.3)
lw = Inches(7.0)
lh = Inches(1.0)
ly_start = Inches(1.2)
for i, (col, name, desc) in enumerate(layers):
    ly = ly_start + i * (lh + Inches(0.15))
    add_rect(s, lx, ly, lw, lh, fill=col)
    add_text(s, name, lx + Inches(0.15), ly + Inches(0.05), Inches(2.5), Inches(0.45),
             size=Pt(15), bold=True, color=RGBColor(0x0D, 0x11, 0x17))
    add_text(s, desc, lx + Inches(0.15), ly + Inches(0.5), lw - Inches(0.3), Inches(0.45),
             size=Pt(12), color=RGBColor(0x0D, 0x11, 0x17))
    if i < 3:
        add_text(s, "▼", lx + lw/2 - Inches(0.15), ly + lh, Inches(0.3), Inches(0.15),
                 size=Pt(12), color=WHITE, align=PP_ALIGN.CENTER)

bx = Inches(7.7)
add_text(s, "Key Design Properties", bx, Inches(1.2), Inches(5.3), Inches(0.5),
         size=Pt(15), bold=True, color=BLUE)
bullets = [
    "Compiler enforces layer boundaries — upward dependencies cause build failure",
    "New graph tab: add ≤ 3 files in Presentation layer only",
    "New audio source: implement IAudioSource — 2 files max",
]
for i, b in enumerate(bullets):
    add_text(s, "• " + b, bx, Inches(1.8 + i * 1.2), Inches(5.3), Inches(1.1),
             size=Pt(15), color=WHITE)

# ── SLIDE 6: C&C Thread Model ──────────────────────────────────────────────
s = blank_slide(prs)
footer(s, 6)
section_label(s, "02 / Architecture")
slide_title(s, "C&C View: DSP Pipeline Thread Model")
add_text(s, "ADR-001  —  QAS-2 Real-Time Performance · QAS-3 Low Latency",
         Inches(0.4), Inches(1.05), Inches(12), Inches(0.35), size=Pt(14), color=BLUE)

img_path = r"D:\sw_architect\project\2026-3-sw-architect-studio-project\docs\milestone2\final\assets\view3-thread-model-simple.png"
try:
    s.shapes.add_picture(img_path, Inches(0.3), Inches(1.5), Inches(8.5), Inches(3.5))
except Exception:
    add_rect(s, Inches(0.3), Inches(1.5), Inches(8.5), Inches(3.5), fill=DARK_CARD)
    add_text(s, "[view3-thread-model-simple.png]", Inches(0.3), Inches(3.0), Inches(8.5), Inches(0.5),
             size=Pt(14), color=GRAY, align=PP_ALIGN.CENTER)

# Table
tx = Inches(9.1)
headers = ["Metric", "Before", "After"]
rows_t = [
    ("wait_ms avg", "77.4 ms", "0.03 ms (×2,600)"),
    ("Deadline miss", "43%", "0%"),
    ("Backlog", "Present", "None"),
]
cw = [Inches(1.5), Inches(1.2), Inches(1.9)]
cx_list = [tx, tx + cw[0], tx + cw[0] + cw[1]]
row_h2 = Inches(0.55)
hy = Inches(1.5)
for ci, h in enumerate(headers):
    add_rect(s, cx_list[ci], hy, cw[ci], row_h2, fill=RGBColor(0x21, 0x26, 0x2E))
    add_text(s, h, cx_list[ci] + Inches(0.05), hy, cw[ci], row_h2,
             size=Pt(12), bold=True, color=GRAY)
for ri, row in enumerate(rows_t):
    ry = hy + row_h2 + ri * row_h2
    bg = DARK_CARD if ri % 2 == 0 else RGBColor(0x1A, 0x20, 0x28)
    after_color = GREEN if ri < 2 else WHITE
    for ci, val in enumerate(row):
        add_rect(s, cx_list[ci], ry, cw[ci], row_h2, fill=bg)
        col = after_color if ci == 2 else WHITE
        add_text(s, val, cx_list[ci] + Inches(0.05), ry, cw[ci], row_h2,
                 size=Pt(12), color=col)

add_rect(s, Inches(0.3), Inches(5.2), Inches(12.7), Inches(0.8),
         fill=BLUE_DARK, line_color=BLUE)
add_text(s, "DSP thread runs at SCHED_FIFO priority. Ring buffer hands off PCM — no shared state, no mutex on the hot path.",
         Inches(0.45), Inches(5.25), Inches(12.4), Inches(0.7), size=Pt(13), color=WHITE)

# ── SLIDE 7: 4-Layer Module View ──────────────────────────────────────────────
s = blank_slide(prs)
footer(s, 7)
section_label(s, "02 / Architecture")
slide_title(s, "Module View: 4-Layer Allowed-to-Use")
add_text(s, "QAS-4 Extensibility / Modifiability",
         Inches(0.4), Inches(1.05), Inches(12), Inches(0.35), size=Pt(14), color=BLUE)

img_path = r"D:\sw_architect\project\2026-3-sw-architect-studio-project\docs\milestone2\final\assets\view1-layered-module.png"
try:
    s.shapes.add_picture(img_path, Inches(0.3), Inches(1.5), Inches(6.5), Inches(5.0))
except Exception:
    add_rect(s, Inches(0.3), Inches(1.5), Inches(6.5), Inches(5.0), fill=DARK_CARD)
    add_text(s, "[view1-layered-module.png]", Inches(0.3), Inches(3.8), Inches(6.5), Inches(0.5),
             size=Pt(14), color=GRAY, align=PP_ALIGN.CENTER)

rx = Inches(7.1)
stats = ["14 graph tabs implemented", "0 layer violations at compile time", "≤ 3 files per new tab"]
for i, st in enumerate(stats):
    add_rect(s, rx, Inches(1.5 + i * 0.65), Inches(5.9), Inches(0.55), fill=DARK_CARD)
    add_text(s, "• " + st, rx + Inches(0.1), Inches(1.5 + i * 0.65), Inches(5.7), Inches(0.55),
             size=Pt(14), bold=True, color=GREEN)

add_text(s, "W2 S1 (11 tabs): TraceDisplay · VarioDisplay · BeatErrorTrace · BeatNoiseScope · WatchPosition · MultiPosition · LongTerm · EscapementAnalyzer · TimeFreqSpectrogram · WaveformComparison · RateScope",
         rx, Inches(3.65), Inches(5.9), Inches(1.2), size=Pt(11), color=GRAY)
add_text(s, "W2 S2 (+2): FilterScope · SweepScope",
         rx, Inches(4.95), Inches(5.9), Inches(0.4), size=Pt(11), color=GRAY)
add_text(s, "W3 S1 Bonus (+1): RadarChart",
         rx, Inches(5.4), Inches(5.9), Inches(0.4), size=Pt(11), color=ORANGE)

# ── SLIDE 8: IAudioSource + Observer ──────────────────────────────────────────────
s = blank_slide(prs)
footer(s, 8)
section_label(s, "02 / Architecture")
slide_title(s, "IAudioSource & BaseGraphTab Observer")

for i, (title, img_file, caption, adr) in enumerate([
    ("IAudioSource (ADR-005)", "view5-iaudiosource.png", "Adding NetworkWorker: ≤ 2 files changed", "ADR-005"),
    ("BaseGraphTab Observer (ADR-006)", "view2b-observer-module.png", "All 14 tabs receive identical Measurement struct", "ADR-006"),
]):
    px = Inches(0.3) if i == 0 else Inches(6.8)
    pw = Inches(6.2)
    add_text(s, title, px, Inches(1.15), pw, Inches(0.45), size=Pt(15), bold=True, color=BLUE)
    img_path = rf"D:\sw_architect\project\2026-3-sw-architect-studio-project\docs\milestone2\final\assets\{img_file}"
    try:
        s.shapes.add_picture(img_path, px, Inches(1.65), pw, Inches(4.3))
    except Exception:
        add_rect(s, px, Inches(1.65), pw, Inches(4.3), fill=DARK_CARD)
        add_text(s, f"[{img_file}]", px, Inches(3.6), pw, Inches(0.5),
                 size=Pt(13), color=GRAY, align=PP_ALIGN.CENTER)
    add_text(s, caption, px, Inches(6.1), pw, Inches(0.4), size=Pt(12), color=GRAY, italic=True)

# ── SLIDE 9: Experiment Results ──────────────────────────────────────────────
s = blank_slide(prs)
footer(s, 9)
section_label(s, "03 / Experiments & Evaluation")
slide_title(s, "Experiment Results")

exp_rows = [
    ("EXP-01", "Rate/Amplitude match WeiShi No.1000?", "✓ Validated W5", "ADR-003 accepted (96kHz)"),
    ("EXP-02", "RPi sustain 96kHz with 0 dropped blocks?", "✓ 0 dropped at 48/96/192kHz", "SCHED_RR not needed"),
    ("EXP-03", "E2E latency: capture → display?", "✓ 2.05ms avg on RPi", "ADR-001 T2 + Ring Buffer"),
    ("EXP-04", "Files required to add new tab?", "✓ ≤ 3 files, 0 violations", "ADR-006 Observer validated"),
    ("EXP-05", "14-tab FPS on RPi — deadline miss?", "✓ Within budget", "ADR-004 on standby"),
]
e_headers = ["ID", "Question", "Result", "Decision"]
e_cw = [Inches(1.0), Inches(4.5), Inches(3.2), Inches(4.2)]
e_cx = [Inches(0.3), Inches(1.35), Inches(5.9), Inches(9.15)]
e_rh = Inches(0.85)
ey = Inches(1.2)

for ci, h in enumerate(e_headers):
    add_rect(s, e_cx[ci], ey, e_cw[ci], Inches(0.42), fill=RGBColor(0x21, 0x26, 0x2E))
    add_text(s, h, e_cx[ci] + Inches(0.05), ey, e_cw[ci], Inches(0.42),
             size=Pt(13), bold=True, color=GRAY)

for ri, row in enumerate(exp_rows):
    ry = ey + Inches(0.42) + ri * e_rh
    bg = RGBColor(0x0D, 0x1F, 0x12) if ri % 2 == 0 else DARK_CARD
    for ci, val in enumerate(row):
        add_rect(s, e_cx[ci], ry, e_cw[ci], e_rh - Inches(0.05), fill=bg)
        col = GREEN if ci == 2 else WHITE
        add_text(s, val, e_cx[ci] + Inches(0.05), ry + Inches(0.05),
                 e_cw[ci] - Inches(0.1), e_rh - Inches(0.1),
                 size=Pt(12), color=col)

# ── SLIDE 10: Key Numbers ──────────────────────────────────────────────
s = blank_slide(prs)
footer(s, 10)
section_label(s, "03 / Experiments & Evaluation")
slide_title(s, "Key Performance Numbers")

metrics = [
    ("0.03 ms", "DSP queue wait (avg)", "Down from 77.4 ms  ·  ×2,600 improvement", GREEN),
    ("2.05 ms", "End-to-end latency (RPi avg)", "Capture → DSP → GUI update", BLUE),
    ("0", "Dropped audio blocks", "At 96kHz over 10-minute session on RPi", GREEN),
    ("142 / 142", "Unit tests passing", "10 test binaries  ·  0 failures", BLUE),
]
card_w2 = Inches(6.1)
card_h2 = Inches(2.4)
positions = [
    (Inches(0.3), Inches(1.3)),
    (Inches(6.7), Inches(1.3)),
    (Inches(0.3), Inches(3.9)),
    (Inches(6.7), Inches(3.9)),
]
for (mx, my), (num, label, sub, col) in zip(positions, metrics):
    add_rect(s, mx, my, card_w2, card_h2, fill=DARK_CARD, line_color=col, line_width=Pt(2))
    add_text(s, num, mx + Inches(0.2), my + Inches(0.15), card_w2 - Inches(0.3), Inches(0.95),
             size=Pt(42), bold=True, color=col)
    add_text(s, label, mx + Inches(0.2), my + Inches(1.15), card_w2 - Inches(0.3), Inches(0.5),
             size=Pt(15), bold=True, color=WHITE)
    add_text(s, sub, mx + Inches(0.2), my + Inches(1.7), card_w2 - Inches(0.3), Inches(0.5),
             size=Pt(12), color=GRAY)

add_text(s, "Target: E2E < 100ms ✓   Deadline miss: 0% ✓   96kHz sustained ✓",
         Inches(0.3), Inches(6.55), Inches(12.7), Inches(0.4),
         size=Pt(14), bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# ── SLIDE 11: ATAM ──────────────────────────────────────────────
s = blank_slide(prs)
footer(s, 11)
section_label(s, "03 / Experiments & Evaluation")
slide_title(s, "Architecture Evaluation — ATAM")
add_text(s, "Architecture Tradeoff Analysis Method (SEI)",
         Inches(0.4), Inches(1.05), Inches(12), Inches(0.35), size=Pt(14), color=BLUE)

col_defs = [
    ("Sensitivity Points", [
        "SP-1  DSP thread location → wait_ms 0.03ms↔77ms",
        "SP-2  isVisible() guard → bottleneck if removed",
        "SP-3  48kHz halves Beat Error resolution",
        "SP-4  Measurement struct must stay immutable",
    ], ORANGE),
    ("Confirmed Non-Risks", [
        "✓ T2 thread removes queue wait (EXP-03)",
        "✓ 96kHz sustainable on RPi (EXP-02)",
        "✓ Lazy rendering: -85% render calls (EXP-03)",
        "✓ 14 tabs each ≤ 3 files (EXP-04)",
        "✓ 142 unit tests all pass",
    ], GREEN),
    ("Resolved Risk Themes", [
        "✓ Theme 1 RESOLVED",
        "  Rendering-Audio Coupling",
        "  ADR-001+002 → miss 43%→0%",
        "",
        "✓ Theme 2 RESOLVED",
        "  WeiShi accuracy gap",
        "  EXP-01 validated W5",
    ], BLUE),
]
col_x_atam = [Inches(0.3), Inches(4.6), Inches(8.9)]
col_w_atam = Inches(4.0)
for ci, (hdr, items, col) in enumerate(col_defs):
    cx = col_x_atam[ci]
    add_rect(s, cx, Inches(1.5), col_w_atam, Inches(0.45), fill=col)
    add_text(s, hdr, cx + Inches(0.1), Inches(1.5), col_w_atam, Inches(0.45),
             size=Pt(13), bold=True, color=RGBColor(0x0D, 0x11, 0x17))
    add_rect(s, cx, Inches(1.95), col_w_atam, Inches(5.0), fill=DARK_CARD)
    body = "\n".join(items)
    add_text(s, body, cx + Inches(0.1), Inches(2.05), col_w_atam - Inches(0.15), Inches(4.8),
             size=Pt(12), color=WHITE)

# ── SLIDE 12: Use of AI ──────────────────────────────────────────────
s = blank_slide(prs)
footer(s, 12)
section_label(s, "04 / Use of AI")
slide_title(s, "How We Used AI in Development")

ai_items = [
    (GREEN,  "Unit Test Generation",
     "142 tests across 10 binaries. Domain math (Rate/Amplitude/Beat Error) verified without equation-level expertise. NTR-07 domain gap addressed."),
    (BLUE,   "Architecture Design & Review",
     "Threading (T1/T2/T3) and rendering (R1/R2/R3) options evaluated via AI dialogue. ADR rationale stress-tested against tradeoff scenarios."),
    (BLUE,   "Debugging & Code Analysis",
     "isVisible()+processEvents() fix diagnosed with AI. Ring buffer sizing and Qt signal/slot thread safety analysis assisted by AI."),
    (ORANGE, "Limitations & Risks",
     "Tests required manual review for correctness assumptions. No hardware access — RPi performance required physical measurement. Hallucination risk on domain equations."),
]
for i, (col, title, desc) in enumerate(ai_items):
    iy = Inches(1.3 + i * 1.45)
    # circle
    add_rect(s, Inches(0.3), iy + Inches(0.1), Inches(0.55), Inches(0.55), fill=col)
    add_text(s, title, Inches(1.05), iy, Inches(11.8), Inches(0.45),
             size=Pt(15), bold=True, color=col)
    add_text(s, desc, Inches(1.05), iy + Inches(0.45), Inches(11.8), Inches(0.85),
             size=Pt(13), color=WHITE)

# ── SLIDE 13: AI Feature ──────────────────────────────────────────────
s = blank_slide(prs)
footer(s, 13)
section_label(s, "04 / Use of AI")
slide_title(s, "AI Feature: On-Device Watch Diagnosis")

# Problem box
add_rect(s, Inches(0.3), Inches(1.2), Inches(7.8), Inches(1.1), fill=BLUE_DARK, line_color=BLUE)
add_text(s, "Domain experts identify watch condition from Rate, Amplitude, and Beat Error — a skill that takes years. Can we automate preliminary diagnosis?",
         Inches(0.45), Inches(1.25), Inches(7.5), Inches(1.0), size=Pt(13), color=WHITE, italic=True)

add_text(s, "Solution: Rule-based + ML classifier on Raspberry Pi 5 (ONNX, fully offline)",
         Inches(0.3), Inches(2.45), Inches(7.8), Inches(0.45), size=Pt(14), bold=True, color=GREEN)

bullets_ai = [
    "Classifies condition: Good / Needs Adjustment / Worn",
    "Flags anomalies: high Beat Error, amplitude drop, rate drift",
    "Provides human-readable explanation per diagnosis",
    "< 50ms inference on RPi 5 — no network required",
]
for i, b in enumerate(bullets_ai):
    add_text(s, "• " + b, Inches(0.5), Inches(3.0 + i * 0.65), Inches(7.5), Inches(0.6),
             size=Pt(14), color=WHITE)

# Right panel
add_rect(s, Inches(8.5), Inches(1.2), Inches(4.5), Inches(4.5), fill=DARK_CARD, line_color=BLUE)
perf = [
    ("Accuracy", "Validated in live demo"),
    ("Latency", "< 50ms inference on RPi 5"),
    ("Model", "ONNX — no Python runtime"),
    ("Mode", "Fully offline, on-device"),
]
for i, (k, v) in enumerate(perf):
    py = Inches(1.4 + i * 0.85)
    add_text(s, k, Inches(8.65), py, Inches(1.4), Inches(0.4), size=Pt(13), bold=True, color=BLUE)
    add_text(s, v, Inches(10.1), py, Inches(2.8), Inches(0.4), size=Pt(13), color=WHITE)

add_text(s, "→  Live demonstration follows",
         Inches(8.65), Inches(5.0), Inches(4.1), Inches(0.4),
         size=Pt(13), color=GRAY, italic=True)

# ── SLIDE 14: What Went Well ──────────────────────────────────────────────
s = blank_slide(prs)
footer(s, 14)
section_label(s, "05 / Lessons Learned")
slide_title(s, "What Went Well")

well_items = [
    ("Experiment-driven architecture decisions",
     "EXP-03 confirmed 43%→0% improvement the same day ADR-001 was implemented. Short feedback loops prevented wasted design work."),
    ("4-Layer structure eliminated regression risk",
     "Compile-time boundary enforcement. The ≤ 3-file constraint held for all 14 tabs — zero coupling surprises."),
    ("AI unit tests solved the domain knowledge gap",
     "NTR-07 (equation verification requires watch expertise) mitigated: 142 tests protect domain math without manual derivation."),
    ("ADR + QA + Experiment traceability",
     "Every architecture decision links to a QA scenario, experiment, and risk. No undocumented choices."),
]
for i, (title, desc) in enumerate(well_items):
    wy = Inches(1.3 + i * 1.45)
    add_rect(s, Inches(0.3), wy, Inches(0.55), Inches(0.55), fill=GREEN)
    add_text(s, "✓", Inches(0.3), wy, Inches(0.55), Inches(0.55),
             size=Pt(18), bold=True, color=RGBColor(0x0D, 0x11, 0x17), align=PP_ALIGN.CENTER)
    add_text(s, title, Inches(1.05), wy, Inches(11.8), Inches(0.45),
             size=Pt(15), bold=True, color=GREEN)
    add_text(s, desc, Inches(1.05), wy + Inches(0.45), Inches(11.8), Inches(0.85),
             size=Pt(13), color=WHITE)

# ── SLIDE 15: What We'd Do Differently ──────────────────────────────────────────────
s = blank_slide(prs)
footer(s, 15)
section_label(s, "05 / Lessons Learned")
slide_title(s, "What We'd Do Differently")

diff_items = [
    ("Start WeiShi validation earlier",
     "QAS-1 (Measurement Accuracy) is the governing goal — but hardware comparison was scheduled last. Next time: validate the top-priority QA first, before structural work."),
    ("Schedule EXP-05 earlier",
     "14-tab FPS test on RPi deferred to W4. ADR-004 remained unconfirmed until the final week. Plan: run full-load experiments in parallel with implementation."),
    ("Stress-test ring buffer depth",
     "Set conservatively without a soak test. ATAM flagged this as R-2. A 30-minute RPi stress test would have given early confidence in the sizing."),
]
for i, (title, desc) in enumerate(diff_items):
    dy = Inches(1.3 + i * 1.6)
    add_rect(s, Inches(0.3), dy, Inches(0.55), Inches(0.55), fill=ORANGE)
    add_text(s, "!", Inches(0.3), dy, Inches(0.55), Inches(0.55),
             size=Pt(18), bold=True, color=RGBColor(0x0D, 0x11, 0x17), align=PP_ALIGN.CENTER)
    add_text(s, title, Inches(1.05), dy, Inches(11.8), Inches(0.45),
             size=Pt(15), bold=True, color=ORANGE)
    add_text(s, desc, Inches(1.05), dy + Inches(0.45), Inches(11.8), Inches(1.0),
             size=Pt(13), color=WHITE)

add_rect(s, Inches(0.3), Inches(6.2), Inches(12.7), Inches(0.85), fill=BLUE_DARK, line_color=BLUE)
add_text(s, "Key principle: Validate the highest-priority QA (Accuracy) first — not last.",
         Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.65),
         size=Pt(15), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ── SAVE ──────────────────────────────────────────────
out = r"D:\sw_architect\project\2026-3-sw-architect-studio-project\docs\milestone3\TimeGrapher-M3-Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
